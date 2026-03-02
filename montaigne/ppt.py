"""PowerPoint generation from PDF or images, and PPTX input support."""

import re
import shutil
import subprocess
from pathlib import Path
from typing import List, Optional

from .logging import get_logger

logger = get_logger(__name__)

IMAGE_EXTENSIONS = {".jpeg", ".jpg", ".png", ".gif", ".webp", ".bmp", ".tiff"}


# ---------------------------------------------------------------------------
# PPTX input support: extract pages/notes from .pptx files
# ---------------------------------------------------------------------------


def check_libreoffice() -> bool:
    """Return True if LibreOffice (soffice) is available on PATH."""
    for cmd in ("soffice", "libreoffice"):
        if shutil.which(cmd):
            return True
    return False


def extract_pptx_pages(
    pptx_path: Path,
    output_dir: Optional[Path] = None,
    dpi: int = 150,
    add_branding: bool = True,
    logo_path: Optional[Path] = None,
) -> List[Path]:
    """
    Extract slide images from a PPTX file.

    Strategy 1 (preferred): If LibreOffice is available, convert PPTX → PDF
    then delegate to ``extract_pdf_pages()``.

    Strategy 2 (fallback): Use python-pptx to pull the largest embedded image
    from each slide (works well for NotebookLM exports where each slide is a
    full-bleed image).

    Args:
        pptx_path: Path to the .pptx file
        output_dir: Directory for output images (default: {stem}_images/)
        dpi: Resolution for PDF-based extraction (default: 150)
        add_branding: If True, add montaigne.cc logo overlay
        logo_path: Optional path to logo image

    Returns:
        Sorted list of extracted image file paths
    """
    from pptx import Presentation as PptxPresentation

    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    if output_dir is None:
        output_dir = pptx_path.parent / f"{pptx_path.stem}_images"
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    # --- Strategy 1: LibreOffice → PDF → extract_pdf_pages ---
    if check_libreoffice():
        logger.info("Converting PPTX to PDF via LibreOffice...")
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        import tempfile

        with tempfile.TemporaryDirectory(prefix="montaigne_lo_") as tmp:
            subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    tmp,
                    str(pptx_path),
                ],
                capture_output=True,
                check=True,
            )
            pdf_candidates = list(Path(tmp).glob("*.pdf"))
            if pdf_candidates:
                from .pdf import extract_pdf_pages

                return extract_pdf_pages(
                    pdf_candidates[0],
                    output_dir=output_dir,
                    dpi=dpi,
                    add_branding=add_branding,
                    logo_path=logo_path,
                )

        logger.warning("LibreOffice conversion produced no PDF, falling back to image extraction")

    # --- Strategy 2: extract embedded images via python-pptx ---
    logger.info("Extracting embedded images from PPTX slides...")
    prs = PptxPresentation(str(pptx_path))
    extracted: List[Path] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        # Collect all images in this slide and pick the largest
        best_blob: Optional[bytes] = None
        best_size = 0
        best_ext = ".png"

        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                blob = shape.image.blob
                if len(blob) > best_size:
                    best_blob = blob
                    best_size = len(blob)
                    content_type = shape.image.content_type  # e.g. "image/png"
                    ext = {
                        "image/png": ".png",
                        "image/jpeg": ".jpg",
                        "image/gif": ".gif",
                        "image/bmp": ".bmp",
                        "image/tiff": ".tiff",
                    }.get(content_type, ".png")
                    best_ext = ext

        if best_blob:
            out_path = output_dir / f"page_{slide_idx:03d}{best_ext}"
            out_path.write_bytes(best_blob)
            extracted.append(out_path)
            logger.info("  Extracted slide %d image (%d KB)", slide_idx, best_size // 1024)
        else:
            logger.warning("  Slide %d: no embedded image found, skipping", slide_idx)

    if not extracted:
        raise RuntimeError(
            f"Could not extract any images from {pptx_path.name}. "
            "Install LibreOffice for full rendering support: brew install --cask libreoffice"
        )

    logger.info("Extracted %d slide images to %s/", len(extracted), output_dir)
    return extracted


def extract_pptx_notes(pptx_path: Path) -> List[str]:
    """
    Extract speaker notes from each slide in a PPTX file.

    Args:
        pptx_path: Path to the .pptx file

    Returns:
        List of note strings, one per slide (empty string if no notes)
    """
    from pptx import Presentation as PptxPresentation

    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    prs = PptxPresentation(str(pptx_path))
    notes: List[str] = []

    for slide in prs.slides:
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text.strip()
            notes.append(text)
        else:
            notes.append("")

    return notes


def notes_to_voiceover_script(
    notes: List[str], output_path: Path, title: str = "Presentation"
) -> Path:
    """
    Format slide notes into the standard voiceover markdown script format.

    Args:
        notes: List of note strings (one per slide)
        output_path: Path where the markdown file will be written
        title: Presentation title for the header

    Returns:
        Path to the written script file
    """
    output_path = Path(output_path)
    lines = [
        f"# {title}",
        "## Voice-Over Script",
        "",
        f"**Total slides:** {len(notes)}",
        "",
        "---",
        "",
    ]

    for i, note_text in enumerate(notes, start=1):
        text = note_text.strip() if note_text else "[No notes for this slide]"
        lines.extend(
            [
                f"## SLIDE {i}: Slide {i}",
                "**Duration:** 30 seconds",
                "",
                "### Voice-Over:",
                "",
                text,
                "",
                "---",
                "",
            ]
        )

    lines.append("*Script generated from PPTX slide notes by Montaigne*")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text("\n".join(lines), encoding="utf-8")
    logger.info("Generated voiceover script from notes: %s", output_path)
    return output_path


def parse_script_to_slides(script_path: Path) -> List[str]:
    """
    Parse a voiceover script markdown file and extract text for each slide.

    Expects format like:
        ## SLIDE 1: Title
        ...script text...
        ---
        ## SLIDE 2: Title
        ...

    Args:
        script_path: Path to the markdown script file

    Returns:
        List of script texts, one per slide (index 0 = slide 1)
    """
    script_path = Path(script_path)
    if not script_path.exists():
        raise FileNotFoundError(f"Script not found: {script_path}")

    with open(script_path, "r", encoding="utf-8") as f:
        content = f.read()

    # Split by slide headers (## SLIDE N: or ## SLIDE N —)
    slide_pattern = r"##\s+SLIDE\s+\d+[:\s—–-]"
    parts = re.split(slide_pattern, content, flags=re.IGNORECASE)

    # First part is header content before first slide, skip it
    slide_texts = []
    for part in parts[1:]:
        # Remove duration markers and separators
        text = re.sub(r"\*\*\[Duration:[^\]]*\]\*\*", "", part)
        text = re.sub(r"^---\s*$", "", text, flags=re.MULTILINE)
        # Clean up extra whitespace
        text = "\n".join(line.strip() for line in text.strip().split("\n"))
        text = re.sub(r"\n{3,}", "\n\n", text)
        slide_texts.append(text.strip())

    return slide_texts


def images_to_pptx(
    images: List[Path], output_path: Path, notes: Optional[List[str]] = None
) -> Path:
    """
    Create a PowerPoint presentation from a list of images.

    Args:
        images: List of image file paths
        output_path: Path for output .pptx file
        notes: Optional list of notes text, one per slide

    Returns:
        Path to the created PowerPoint file
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    for i, image_path in enumerate(images):
        image_path = Path(image_path)
        if not image_path.exists():
            logger.warning("Image not found, skipping: %s", image_path)
            continue

        slide = prs.slides.add_slide(blank_layout)

        # Add image to fill the slide
        # Calculate dimensions to fit while maintaining aspect ratio
        from PIL import Image

        with Image.open(image_path) as img:
            img_width, img_height = img.size

        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Scale to fit slide
        width_ratio = slide_width / Inches(img_width / 96)  # assuming 96 DPI
        height_ratio = slide_height / Inches(img_height / 96)
        scale = min(width_ratio, height_ratio, 1.0)  # Don't scale up

        pic_width = Inches(img_width / 96) * scale
        pic_height = Inches(img_height / 96) * scale

        # Center the image
        left = (slide_width - pic_width) / 2
        top = (slide_height - pic_height) / 2

        slide.shapes.add_picture(str(image_path), left, top, pic_width, pic_height)

        # Add notes if provided
        if notes and i < len(notes) and notes[i]:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes[i]

        logger.info("Added slide %d: %s", i + 1, image_path.name)

    prs.save(output_path)
    return output_path


def pdf_to_pptx(
    pdf_path: Path,
    output_path: Optional[Path] = None,
    script_path: Optional[Path] = None,
    dpi: int = 150,
    keep_images: bool = False,
) -> Path:
    """
    Convert a PDF to a PowerPoint presentation.

    Each page of the PDF becomes a slide with the page as an image.

    Args:
        pdf_path: Path to the PDF file
        output_path: Path for output .pptx file (default: {pdf_stem}.pptx)
        script_path: Optional path to voiceover script for slide notes
        dpi: Resolution for PDF extraction (default: 150)
        keep_images: If True, keep extracted images; if False, delete them

    Returns:
        Path to the created PowerPoint file
    """
    from .pdf import extract_pdf_pages
    import tempfile
    import shutil

    pdf_path = Path(pdf_path)
    if not pdf_path.exists():
        raise FileNotFoundError(f"PDF not found: {pdf_path}")

    if output_path is None:
        output_path = pdf_path.parent / f"{pdf_path.stem}.pptx"
    output_path = Path(output_path)

    logger.info("Converting PDF to PowerPoint: %s", pdf_path.name)

    # Extract PDF pages to temporary directory or keep
    if keep_images:
        images_dir = pdf_path.parent / f"{pdf_path.stem}_images"
    else:
        images_dir = Path(tempfile.mkdtemp(prefix="montaigne_pdf_"))

    try:
        images = extract_pdf_pages(pdf_path, output_dir=images_dir, dpi=dpi)

        # Parse script if provided
        notes = None
        if script_path:
            script_path = Path(script_path)
            logger.info("Parsing script for notes: %s", script_path.name)
            notes = parse_script_to_slides(script_path)
            if len(notes) != len(images):
                logger.warning("Script has %d slides but PDF has %d pages", len(notes), len(images))

        # Create PowerPoint
        logger.info("Creating PowerPoint with %d slides...", len(images))
        images_to_pptx(images, output_path, notes=notes)

    finally:
        # Clean up temporary images if not keeping
        if not keep_images and images_dir.exists():
            shutil.rmtree(images_dir)

    logger.info("Created: %s", output_path)
    return output_path


def folder_to_pptx(
    folder_path: Path, output_path: Optional[Path] = None, script_path: Optional[Path] = None
) -> Path:
    """
    Convert a folder of images to a PowerPoint presentation.

    Images are sorted alphabetically/numerically and each becomes a slide.

    Args:
        folder_path: Path to folder containing images
        output_path: Path for output .pptx file (default: {folder_name}.pptx)
        script_path: Optional path to voiceover script for slide notes

    Returns:
        Path to the created PowerPoint file
    """
    folder_path = Path(folder_path)
    if not folder_path.exists():
        raise FileNotFoundError(f"Folder not found: {folder_path}")
    if not folder_path.is_dir():
        raise ValueError(f"Not a directory: {folder_path}")

    # Find all images in folder
    images = sorted([f for f in folder_path.iterdir() if f.suffix.lower() in IMAGE_EXTENSIONS])

    if not images:
        raise ValueError(f"No images found in {folder_path}")

    if output_path is None:
        output_path = folder_path.parent / f"{folder_path.name}.pptx"
    output_path = Path(output_path)

    logger.info("Creating PowerPoint from %d images in: %s", len(images), folder_path.name)

    # Parse script if provided
    notes = None
    if script_path:
        script_path = Path(script_path)
        logger.info("Parsing script for notes: %s", script_path.name)
        notes = parse_script_to_slides(script_path)
        if len(notes) != len(images):
            logger.warning("Script has %d slides but folder has %d images", len(notes), len(images))

    # Create PowerPoint
    logger.info("Creating PowerPoint with %d slides...", len(images))
    images_to_pptx(images, output_path, notes=notes)

    logger.info("Created: %s", output_path)
    return output_path


def create_pptx(
    input_path: Path,
    output_path: Optional[Path] = None,
    script_path: Optional[Path] = None,
    dpi: int = 150,
    keep_images: bool = False,
) -> Path:
    """
    Create a PowerPoint presentation from PDF or image folder.

    This is the main entry point that auto-detects input type.

    Args:
        input_path: Path to PDF file or folder of images
        output_path: Path for output .pptx file
        script_path: Optional path to voiceover script for slide notes
        dpi: Resolution for PDF extraction (default: 150)
        keep_images: If True and input is PDF, keep extracted images

    Returns:
        Path to the created PowerPoint file
    """
    input_path = Path(input_path)

    if input_path.suffix.lower() == ".pdf":
        return pdf_to_pptx(
            input_path,
            output_path=output_path,
            script_path=script_path,
            dpi=dpi,
            keep_images=keep_images,
        )
    elif input_path.is_dir():
        return folder_to_pptx(input_path, output_path=output_path, script_path=script_path)
    else:
        raise ValueError(f"Input must be a PDF file or folder of images: {input_path}")
