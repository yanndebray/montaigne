"""Tests for PPTX input support — notes extraction, image extraction, voiceover script generation."""

import io
import pytest
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from montaigne.ppt import (
    check_libreoffice,
    extract_pptx_notes,
    extract_pptx_pages,
    notes_to_voiceover_script,
    pptx_has_notes,
)


# ---------------------------------------------------------------------------
# Helpers to create test PPTX files
# ---------------------------------------------------------------------------


def _make_pptx_with_notes(path: Path, notes: list[str]) -> Path:
    """Create a PPTX file where each slide has the given note text."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for note_text in notes:
        slide = prs.slides.add_slide(blank_layout)
        if note_text:
            slide.notes_slide.notes_text_frame.text = note_text
    prs.save(str(path))
    return path


def _make_pptx_with_images(path: Path, num_slides: int = 2) -> Path:
    """Create a PPTX where each slide has an embedded PNG image."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for i in range(num_slides):
        slide = prs.slides.add_slide(blank_layout)
        # Create a small in-memory PNG
        img = Image.new("RGB", (200, 150), color=(50 * i, 100, 200))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        slide.shapes.add_picture(buf, Inches(0), Inches(0), Inches(10), Inches(7.5))
    prs.save(str(path))
    return path


# ---------------------------------------------------------------------------
# Tests: check_libreoffice
# ---------------------------------------------------------------------------


class TestCheckLibreOffice:
    def test_returns_true_when_soffice_on_path(self):
        with patch("montaigne.ppt.shutil.which", side_effect=lambda cmd: "/usr/bin/soffice" if cmd == "soffice" else None):
            assert check_libreoffice() is True

    def test_returns_true_when_libreoffice_on_path(self):
        with patch("montaigne.ppt.shutil.which", side_effect=lambda cmd: "/usr/bin/libreoffice" if cmd == "libreoffice" else None):
            assert check_libreoffice() is True

    def test_returns_false_when_neither_available(self):
        with patch("montaigne.ppt.shutil.which", return_value=None):
            assert check_libreoffice() is False


# ---------------------------------------------------------------------------
# Tests: extract_pptx_notes
# ---------------------------------------------------------------------------


class TestExtractPptxNotes:
    def test_extracts_notes_from_slides(self, temp_dir):
        notes = ["Welcome to slide one.", "Key details here.", "Thank you."]
        pptx_path = _make_pptx_with_notes(temp_dir / "deck.pptx", notes)

        result = extract_pptx_notes(pptx_path)

        assert len(result) == 3
        assert result[0] == "Welcome to slide one."
        assert result[1] == "Key details here."
        assert result[2] == "Thank you."

    def test_empty_notes_returns_empty_strings(self, temp_dir):
        pptx_path = _make_pptx_with_notes(temp_dir / "empty.pptx", ["", "", ""])

        result = extract_pptx_notes(pptx_path)

        assert len(result) == 3
        assert all(n == "" for n in result)

    def test_mixed_notes_and_empty(self, temp_dir):
        pptx_path = _make_pptx_with_notes(temp_dir / "mixed.pptx", ["Has notes", "", "Also notes"])

        result = extract_pptx_notes(pptx_path)

        assert result[0] == "Has notes"
        assert result[1] == ""
        assert result[2] == "Also notes"

    def test_nonexistent_file_raises(self, temp_dir):
        with pytest.raises(FileNotFoundError):
            extract_pptx_notes(temp_dir / "nope.pptx")


# ---------------------------------------------------------------------------
# Tests: pptx_has_notes
# ---------------------------------------------------------------------------


class TestPptxHasNotes:
    def test_true_when_notes_present(self, temp_dir):
        pptx_path = _make_pptx_with_notes(temp_dir / "with_notes.pptx", ["Hello", "", "World"])
        assert pptx_has_notes(pptx_path) is True

    def test_false_when_no_notes(self, temp_dir):
        pptx_path = _make_pptx_with_notes(temp_dir / "empty.pptx", ["", "", ""])
        assert pptx_has_notes(pptx_path) is False

    def test_false_when_only_whitespace(self, temp_dir):
        pptx_path = _make_pptx_with_notes(temp_dir / "ws.pptx", ["  ", "\n", ""])
        assert pptx_has_notes(pptx_path) is False


# ---------------------------------------------------------------------------
# Tests: notes_to_voiceover_script
# ---------------------------------------------------------------------------


class TestNotesToVoiceoverScript:
    def test_generates_standard_format(self, temp_dir):
        notes = ["Intro text.", "Body text.", "Closing text."]
        output = temp_dir / "voiceover.md"

        result = notes_to_voiceover_script(notes, output, title="My Deck")

        assert result == output
        assert output.exists()

        content = output.read_text(encoding="utf-8")
        assert "# My Deck" in content
        assert "## SLIDE 1:" in content
        assert "## SLIDE 2:" in content
        assert "## SLIDE 3:" in content
        assert "Intro text." in content
        assert "Body text." in content
        assert "Closing text." in content

    def test_empty_note_gets_placeholder(self, temp_dir):
        notes = ["", "Has text"]
        output = temp_dir / "voiceover.md"

        notes_to_voiceover_script(notes, output)

        content = output.read_text(encoding="utf-8")
        assert "[No notes for this slide]" in content
        assert "Has text" in content

    def test_creates_parent_directories(self, temp_dir):
        output = temp_dir / "sub" / "dir" / "voiceover.md"

        notes_to_voiceover_script(["Hello"], output)

        assert output.exists()


# ---------------------------------------------------------------------------
# Tests: extract_pptx_pages (image extraction fallback)
# ---------------------------------------------------------------------------


class TestExtractPptxPages:
    def test_extracts_images_from_slides(self, temp_dir):
        """With LibreOffice unavailable, extract embedded images via python-pptx."""
        pptx_path = _make_pptx_with_images(temp_dir / "slides.pptx", num_slides=3)
        output_dir = temp_dir / "output"

        with patch("montaigne.ppt.check_libreoffice", return_value=False):
            result = extract_pptx_pages(pptx_path, output_dir=output_dir)

        assert len(result) == 3
        assert all(p.exists() for p in result)
        assert result[0].name == "page_001.png"
        assert result[1].name == "page_002.png"
        assert result[2].name == "page_003.png"

    def test_nonexistent_file_raises(self, temp_dir):
        with pytest.raises(FileNotFoundError):
            extract_pptx_pages(temp_dir / "missing.pptx")

    def test_no_images_and_no_libreoffice_raises(self, temp_dir):
        """A PPTX with no embedded images and no LibreOffice should raise RuntimeError."""
        # Create PPTX with no images (just blank slides)
        pptx_path = _make_pptx_with_notes(temp_dir / "blank.pptx", ["note"])

        with patch("montaigne.ppt.check_libreoffice", return_value=False):
            with pytest.raises(RuntimeError, match="Could not extract any images"):
                extract_pptx_pages(pptx_path)

    def test_default_output_dir(self, temp_dir):
        """Default output directory should be {stem}_images/."""
        pptx_path = _make_pptx_with_images(temp_dir / "deck.pptx", num_slides=1)

        with patch("montaigne.ppt.check_libreoffice", return_value=False):
            result = extract_pptx_pages(pptx_path)

        expected_dir = temp_dir / "deck_images"
        assert expected_dir.exists()
        assert result[0].parent == expected_dir
